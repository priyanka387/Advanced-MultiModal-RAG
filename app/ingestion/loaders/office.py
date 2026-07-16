import os
import logfire

from docx import Document
from pptx import Presentation


def parse_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    doc = Document(file_path)

    text = []

    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                text.append(" | ".join(values))

    return "\n".join(text)


def parse_pptx(file_path: str) -> str:
    """Extract text from PowerPoint."""
    prs = Presentation(file_path)

    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)

    return "\n".join(text)


def parse_office(file_path: str) -> str:

    ext = os.path.splitext(file_path)[1].lower()

    with logfire.span("Office Parsing", filename=file_path):

        if ext == ".docx":
            text = parse_docx(file_path)

        elif ext == ".pptx":
            text = parse_pptx(file_path)

        else:
            raise ValueError(f"Unsupported Office file: {ext}")

        logfire.info(f"Extracted {len(text)} characters")

        return text